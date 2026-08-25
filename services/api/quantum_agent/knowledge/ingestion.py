"""Deterministic, provenance-preserving course-material ingestion.

This module deliberately stops at source parsing.  It does not call an LLM,
write to a database, or approve graph knowledge.  Every output can therefore
be reproduced from the source bytes and reviewed before graph extraction.

The parsers retain two representations:

* ``SourceUnit.raw_text`` is the exact text returned by the format parser for
  one page, slide, paragraph, row, or line range.
* ``SourceChunk.exact_text`` is an exact slice of the retained unit content.
  Repeated PDF/PPTX marginalia may be excluded, but the excluded text remains
  in the unit's blocks and ``removed_marginalia`` field.

PDF page labels are populated only when the PDF contains a PageLabels rule.
The one-based physical page number is always stored separately.
"""

# Course-authored Chinese punctuation and numerals are intentionally exact.
# ruff: noqa: RUF001

from __future__ import annotations

import hashlib
import json
import math
import mimetypes
import re
import unicodedata
from collections import defaultdict
from collections.abc import Iterable, Iterator, Mapping, Sequence
from dataclasses import dataclass, field
from datetime import date, datetime, time
from enum import StrEnum
from pathlib import Path
from typing import Any, Literal

from pydantic import BaseModel, ConfigDict, Field, model_validator

# Bumped when paragraph/line blocks became first-class retained text.  Persisted
# immutable versions use this value to reject silent parser-behavior drift.
PARSER_VERSION = "1.1.0"


class IngestionError(RuntimeError):
    """Base exception for deterministic ingestion failures."""


class UnsupportedDocumentError(IngestionError):
    """Raised when no parser is registered for a file suffix."""


class MissingParserDependency(IngestionError):
    """Raised when an optional parser dependency is unavailable."""


class LocatorType(StrEnum):
    PAGE = "page"
    SLIDE = "slide"
    PARAGRAPH = "paragraph"
    SHEET_ROW = "sheet_row"
    LINE = "line"


class IngestionStatus(StrEnum):
    READY = "READY"
    REVIEW_REQUIRED = "REVIEW_REQUIRED"
    OCR_REQUIRED = "OCR_REQUIRED"


class SourceLocator(BaseModel):
    """A format-native, non-invented source location."""

    model_config = ConfigDict(extra="forbid", frozen=True)

    locator_type: LocatorType
    start: int | str
    end: int | str
    physical_page: int | None = None
    page_label: str | None = None
    sheet_name: str | None = None

    @model_validator(mode="after")
    def validate_location(self) -> SourceLocator:
        if self.page_label == "":
            raise ValueError("empty page labels must be represented as null")
        if self.locator_type == LocatorType.PAGE:
            if self.physical_page is None or self.physical_page < 1:
                raise ValueError("page locators require a one-based physical_page")
        elif self.physical_page is not None or self.page_label is not None:
            raise ValueError("physical_page and page_label are valid only for page locators")
        if self.locator_type == LocatorType.SHEET_ROW and not self.sheet_name:
            raise ValueError("sheet-row locators require sheet_name")
        if self.locator_type != LocatorType.SHEET_ROW and self.sheet_name is not None:
            raise ValueError("sheet_name is valid only for sheet-row locators")
        return self


class TextBlock(BaseModel):
    """One parser-returned text or image block with its native bounding box."""

    model_config = ConfigDict(extra="forbid", frozen=True)

    ordinal: int
    exact_text: str
    bbox: tuple[float, float, float, float] | None = None
    kind: str = "text"
    region: Literal["top", "body", "bottom"] = "body"
    removed_as_repeated_marginalia: bool = False
    metadata: Mapping[str, str | int | float | bool | None] = Field(default_factory=dict)


class SourceUnit(BaseModel):
    """The smallest source boundary across which chunks must never cross."""

    model_config = ConfigDict(extra="forbid", frozen=True)

    id: str
    document_version_id: str
    locator: SourceLocator
    section_path: tuple[str, ...] = ()
    raw_text: str
    content_text: str
    raw_checksum: str
    blocks: tuple[TextBlock, ...] = ()
    removed_marginalia: tuple[str, ...] = ()
    status: IngestionStatus = IngestionStatus.READY
    flags: tuple[str, ...] = ()
    render_required: bool = False
    metadata: Mapping[str, str | int | float | bool | None] = Field(default_factory=dict)


class SourceChunk(BaseModel):
    """A deterministic, auditable slice used as extraction/retrieval evidence."""

    model_config = ConfigDict(extra="forbid", frozen=True)

    id: str
    document_id: str
    document_version_id: str
    source_unit_id: str
    ordinal: int
    exact_text: str
    checksum: str
    section_path: tuple[str, ...] = ()
    locator: SourceLocator
    content_char_start: int
    content_char_end: int
    evidence_snippet: str
    evidence_char_start: int
    evidence_char_end: int
    evidence_snippet_basis: Literal["exact_text"] = "exact_text"
    status: IngestionStatus = IngestionStatus.READY
    flags: tuple[str, ...] = ()

    @model_validator(mode="after")
    def validate_evidence(self) -> SourceChunk:
        if self.content_char_end < self.content_char_start:
            raise ValueError("invalid content character range")
        if not 0 <= self.evidence_char_start <= self.evidence_char_end <= len(self.exact_text):
            raise ValueError("invalid evidence character range")
        expected = self.exact_text[self.evidence_char_start : self.evidence_char_end]
        if self.evidence_snippet != expected:
            raise ValueError("evidence_snippet must be an exact substring of exact_text")
        if self.checksum != sha256_text(self.exact_text):
            raise ValueError("chunk checksum does not match exact_text")
        return self


class IngestedDocument(BaseModel):
    """Pure parse result; no graph knowledge has been approved at this stage."""

    model_config = ConfigDict(extra="forbid", frozen=True)

    document_id: str
    document_version_id: str
    source_path: str
    filename: str
    media_type: str
    sha256: str
    byte_size: int
    parser_name: str
    parser_version: str = PARSER_VERSION
    units: tuple[SourceUnit, ...]
    chunks: tuple[SourceChunk, ...]
    metadata: Mapping[str, str | int | float | bool | None] = Field(default_factory=dict)


class IngestionConfig(BaseModel):
    """Bounds for format-native chunking and conservative quality detection."""

    model_config = ConfigDict(extra="forbid", frozen=True)

    max_chunk_chars: int = Field(default=1_800, ge=128)
    chunk_overlap_chars: int = Field(default=160, ge=0)
    evidence_snippet_chars: int = Field(default=360, ge=32)
    pdf_low_text_characters: int = Field(default=40, ge=0)
    slide_low_text_characters: int = Field(default=20, ge=0)
    repeated_margin_min_units: int = Field(default=3, ge=3)
    repeated_margin_ratio: float = Field(default=0.60, ge=0.5, le=1.0)

    @model_validator(mode="after")
    def validate_chunk_overlap(self) -> IngestionConfig:
        if self.chunk_overlap_chars >= self.max_chunk_chars:
            raise ValueError("chunk_overlap_chars must be smaller than max_chunk_chars")
        return self


@dataclass(slots=True)
class _DraftBlock:
    ordinal: int
    exact_text: str
    bbox: tuple[float, float, float, float] | None = None
    kind: str = "text"
    region: Literal["top", "body", "bottom"] = "body"
    metadata: dict[str, str | int | float | bool | None] = field(default_factory=dict)
    removed: bool = False


@dataclass(slots=True)
class _UnitDraft:
    locator: SourceLocator
    section_path: tuple[str, ...]
    raw_text: str
    blocks: list[_DraftBlock]
    joiner: str = ""
    has_images: bool = False
    low_text_threshold: int = 0
    force_review_flags: set[str] = field(default_factory=set)
    metadata: dict[str, str | int | float | bool | None] = field(default_factory=dict)


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def sha256_text(text: str) -> str:
    return sha256_bytes(text.encode("utf-8"))


def _sanitize_extracted_text(text: str) -> str:
    """Remove bytes PostgreSQL UTF8 text columns reject without changing content.

    Some PDF embedded fonts emit NUL (0x00) and other C0 control bytes outside
    tab/newline/carriage-return.  SQLite silently stores them but
    ``asyncpg`` raises ``CharacterNotInRepertoireError``; stripped here so the
    authoritative chunks persist identically on either backend.
    """

    return "".join(
        character for character in text
        if character not in "\x00\x01\x02\x03\x04\x05\x06\x07\x08\x0b\x0c\x0e"
        "\x0f\x10\x11\x12\x13\x14\x15\x16\x17\x18\x19\x1a\x1b\x1c\x1d\x1e\x1f\x7f"
    )


def _stable_id(prefix: str, *parts: object) -> str:
    payload = json.dumps(
        parts,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
        default=str,
    ).encode("utf-8")
    return f"{prefix}_{hashlib.sha256(payload).hexdigest()[:24]}"


def _logical_source_name(path: Path, source_name: str | None) -> str:
    value = source_name if source_name is not None else path.name
    return unicodedata.normalize("NFC", value.replace("\\", "/"))


def _media_type(path: Path) -> str:
    explicit = {
        ".md": "text/markdown",
        ".txt": "text/plain",
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }
    return explicit.get(
        path.suffix.casefold(),
        mimetypes.guess_type(path.name)[0] or "application/octet-stream",
    )


def _meaningful_character_count(text: str) -> int:
    return sum(character.isalnum() for character in text)


def _margin_signature(text: str) -> str | None:
    normalized = " ".join(text.split()).casefold()
    if not normalized or len(normalized) > 160 or normalized.count(" ") > 24:
        return None
    if re.fullmatch(r"(?:page\s*)?\d+(?:\s*(?:/|of)\s*\d+)?", normalized):
        return "<varying-page-number>"
    if re.fullmatch(r"[ivxlcdm]{1,10}", normalized):
        return "<varying-roman-page-number>"
    return normalized


def _remove_repeated_marginalia(drafts: Sequence[_UnitDraft], config: IngestionConfig) -> None:
    if len(drafts) < config.repeated_margin_min_units:
        return
    occurrences: dict[tuple[str, str], list[tuple[int, _DraftBlock]]] = defaultdict(list)
    for unit_index, draft in enumerate(drafts):
        for block in draft.blocks:
            if block.kind != "text" or block.region == "body":
                continue
            signature = _margin_signature(block.exact_text)
            if signature is not None:
                occurrences[(block.region, signature)].append((unit_index, block))

    threshold = max(
        config.repeated_margin_min_units,
        math.ceil(len(drafts) * config.repeated_margin_ratio),
    )
    for matched in occurrences.values():
        if len({unit_index for unit_index, _ in matched}) < threshold:
            continue
        for _, block in matched:
            block.removed = True


def _preferred_break(text: str, start: int, hard_end: int) -> int:
    if hard_end >= len(text):
        return len(text)
    lower_bound = start + max(1, (hard_end - start) // 2)
    for separator in ("\n\n", "\n", "。", "！", "？", ". ", "; ", " "):
        position = text.rfind(separator, lower_bound, hard_end)
        if position >= lower_bound:
            return position + len(separator)
    return hard_end


def _exact_spans(text: str, config: IngestionConfig) -> list[tuple[int, int]]:
    if not text:
        return [(0, 0)]
    spans: list[tuple[int, int]] = []
    start = 0
    while start < len(text):
        hard_end = min(len(text), start + config.max_chunk_chars)
        end = _preferred_break(text, start, hard_end)
        if end <= start:
            end = hard_end
        spans.append((start, end))
        if end == len(text):
            break
        next_start = end - config.chunk_overlap_chars
        start = max(start + 1, next_start)
    return spans


def _snippet(text: str, limit: int) -> tuple[str, int, int]:
    if not text:
        return "", 0, 0
    first = next((index for index, character in enumerate(text) if not character.isspace()), 0)
    end = min(len(text), first + limit)
    return text[first:end], first, end


def _chunk_locator(locator: SourceLocator, content: str, start: int, end: int) -> SourceLocator:
    if locator.locator_type != LocatorType.LINE:
        return locator
    base = int(locator.start)
    start_line = base + content[:start].count("\n")
    if end <= start:
        end_line = start_line
    else:
        end_line = base + content[: max(start, end - 1)].count("\n")
    return SourceLocator(locator_type=LocatorType.LINE, start=start_line, end=end_line)


def _build_document(
    path: Path,
    *,
    parser_name: str,
    drafts: Sequence[_UnitDraft],
    config: IngestionConfig,
    source_name: str | None,
    remove_repeated_marginalia: bool = False,
    metadata: Mapping[str, str | int | float | bool | None] | None = None,
) -> IngestedDocument:
    source_bytes = path.read_bytes()
    source_sha256 = sha256_bytes(source_bytes)
    logical_name = _logical_source_name(path, source_name)
    document_id = _stable_id("doc", logical_name)
    version_id = _stable_id("docv", document_id, source_sha256, PARSER_VERSION)

    mutable_drafts = list(drafts)
    if remove_repeated_marginalia:
        _remove_repeated_marginalia(mutable_drafts, config)

    units: list[SourceUnit] = []
    chunks: list[SourceChunk] = []
    chunk_ordinal = 0
    for draft in mutable_drafts:
        retained_blocks = [block for block in draft.blocks if not block.removed]
        content_text = draft.joiner.join(
            block.exact_text for block in retained_blocks if block.kind != "image"
        )
        removed_text = tuple(
            block.exact_text for block in draft.blocks if block.removed and block.exact_text
        )
        flags = set(draft.force_review_flags)
        meaningful = _meaningful_character_count(content_text)
        if removed_text:
            flags.add("REPEATED_MARGIN_REMOVED")
        if draft.has_images:
            flags.update(("HAS_IMAGE", "RENDER_REQUIRED"))

        status = IngestionStatus.READY
        if draft.has_images and meaningful == 0:
            status = IngestionStatus.OCR_REQUIRED
            flags.add("IMAGE_ONLY")
        elif draft.has_images and meaningful < draft.low_text_threshold:
            status = IngestionStatus.OCR_REQUIRED
            flags.add("LOW_TEXT")
        elif meaningful < draft.low_text_threshold:
            status = IngestionStatus.REVIEW_REQUIRED
            flags.add("LOW_TEXT")
        elif flags & {"REVIEW_REQUIRED", "PARSER_ANOMALY"}:
            status = IngestionStatus.REVIEW_REQUIRED

        locator_dump = draft.locator.model_dump(mode="json")
        unit_id = _stable_id("unit", version_id, locator_dump, draft.raw_text)
        final_blocks = tuple(
            TextBlock(
                ordinal=block.ordinal,
                exact_text=block.exact_text,
                bbox=block.bbox,
                kind=block.kind,
                region=block.region,
                removed_as_repeated_marginalia=block.removed,
                metadata=block.metadata,
            )
            for block in draft.blocks
        )
        unit = SourceUnit(
            id=unit_id,
            document_version_id=version_id,
            locator=draft.locator,
            section_path=draft.section_path,
            raw_text=draft.raw_text,
            content_text=content_text,
            raw_checksum=sha256_text(draft.raw_text),
            blocks=final_blocks,
            removed_marginalia=removed_text,
            status=status,
            flags=tuple(sorted(flags)),
            render_required=draft.has_images or status == IngestionStatus.OCR_REQUIRED,
            metadata=draft.metadata,
        )
        units.append(unit)

        for content_start, content_end in _exact_spans(content_text, config):
            exact_text = content_text[content_start:content_end]
            evidence, evidence_start, evidence_end = _snippet(
                exact_text, config.evidence_snippet_chars
            )
            chunk_locator = _chunk_locator(draft.locator, content_text, content_start, content_end)
            chunk_id = _stable_id(
                "chunk",
                version_id,
                chunk_locator.model_dump(mode="json"),
                chunk_ordinal,
                exact_text,
            )
            chunks.append(
                SourceChunk(
                    id=chunk_id,
                    document_id=document_id,
                    document_version_id=version_id,
                    source_unit_id=unit_id,
                    ordinal=chunk_ordinal,
                    exact_text=exact_text,
                    checksum=sha256_text(exact_text),
                    section_path=draft.section_path,
                    locator=chunk_locator,
                    content_char_start=content_start,
                    content_char_end=content_end,
                    evidence_snippet=evidence,
                    evidence_char_start=evidence_start,
                    evidence_char_end=evidence_end,
                    status=status,
                    flags=tuple(sorted(flags)),
                )
            )
            chunk_ordinal += 1

    return IngestedDocument(
        document_id=document_id,
        document_version_id=version_id,
        source_path=str(path),
        filename=path.name,
        media_type=_media_type(path),
        sha256=source_sha256,
        byte_size=len(source_bytes),
        parser_name=parser_name,
        units=tuple(units),
        chunks=tuple(chunks),
        metadata=metadata or {},
    )


def _load_pymupdf() -> Any:
    try:
        import pymupdf

        return pymupdf
    except ImportError:
        try:
            import fitz  # type: ignore[import-untyped]

            return fitz
        except ImportError as error:
            raise MissingParserDependency(
                "PDF ingestion requires PyMuPDF (import name 'pymupdf' or 'fitz')"
            ) from error


def _toc_paths(document: Any, page_count: int) -> list[tuple[str, ...]]:
    paths: list[tuple[str, ...]] = [() for _ in range(page_count)]
    try:
        toc = document.get_toc(simple=True)
    except (AttributeError, RuntimeError, ValueError):
        return paths

    events: dict[int, list[tuple[int, str]]] = defaultdict(list)
    for entry in toc:
        if len(entry) < 3:
            continue
        level, title, page_number = entry[:3]
        try:
            level_int = max(1, int(level))
            page_int = int(page_number)
        except (TypeError, ValueError):
            continue
        title_text = _sanitize_extracted_text(str(title)).strip()
        if not title_text or not 1 <= page_int <= page_count:
            continue
        events[page_int - 1].append((level_int, title_text))

    active: list[tuple[int, str]] = []
    for page_index in range(page_count):
        for level, title in events.get(page_index, []):
            active = [
                (old_level, old_title) for old_level, old_title in active if old_level < level
            ]
            active.append((level, title))
        paths[page_index] = tuple(title for _, title in sorted(active, key=lambda item: item[0]))
    return paths


def _pdf_has_explicit_page_labels(document: Any) -> bool:
    """Return true only for an actual PDF PageLabels number tree."""

    try:
        rules = document.get_page_labels()
    except (AttributeError, RuntimeError, ValueError):
        return False
    return bool(rules)


def parse_pdf_document(
    path: str | Path,
    *,
    config: IngestionConfig | None = None,
    source_name: str | None = None,
) -> IngestedDocument:
    """Parse a PDF with one-based physical pages and optional real page labels."""

    pymupdf = _load_pymupdf()
    source_path = Path(path)
    parse_config = config or IngestionConfig()
    drafts: list[_UnitDraft] = []
    try:
        document = pymupdf.open(str(source_path))
    except Exception as error:
        raise IngestionError(f"could not open PDF {source_path.name!r}") from error

    try:
        page_count = int(document.page_count)
        section_paths = _toc_paths(document, page_count)
        has_explicit_labels = _pdf_has_explicit_page_labels(document)
        for page_index in range(page_count):
            page = document.load_page(page_index)
            width = float(page.rect.width)
            height = float(page.rect.height)
            raw_text = _sanitize_extracted_text(page.get_text("text", sort=True))
            blocks: list[_DraftBlock] = []
            raw_blocks = page.get_text("blocks", sort=True)
            for block_ordinal, raw_block in enumerate(raw_blocks):
                if len(raw_block) < 7:
                    continue
                x0, y0, x1, y1 = (round(float(value), 3) for value in raw_block[:4])
                block_type = int(raw_block[6])
                exact_text = (
                    _sanitize_extracted_text(str(raw_block[4])) if block_type == 0 else ""
                )
                if y1 <= height * 0.12:
                    region: Literal["top", "body", "bottom"] = "top"
                elif y0 >= height * 0.88:
                    region = "bottom"
                else:
                    region = "body"
                blocks.append(
                    _DraftBlock(
                        ordinal=block_ordinal,
                        exact_text=exact_text,
                        bbox=(x0, y0, x1, y1),
                        kind="text" if block_type == 0 else "image",
                        region=region,
                        metadata={"pymupdf_block_type": block_type},
                    )
                )

            text_blocks = [block for block in blocks if block.kind == "text"]
            if raw_text and not text_blocks:
                # Defensive fallback: retain parser text if a PyMuPDF build did
                # not return block tuples.  The full-page bbox is auditable.
                blocks.append(
                    _DraftBlock(
                        ordinal=len(blocks),
                        exact_text=raw_text,
                        bbox=(0.0, 0.0, round(width, 3), round(height, 3)),
                        kind="text",
                    )
                )

            try:
                image_count = len(page.get_images(full=True))
            except (AttributeError, RuntimeError, ValueError):
                image_count = sum(block.kind == "image" for block in blocks)
            has_images = image_count > 0 or any(block.kind == "image" for block in blocks)

            page_label: str | None = None
            if has_explicit_labels:
                try:
                    parser_label = page.get_label()
                except (AttributeError, RuntimeError, ValueError):
                    parser_label = None
                if parser_label is not None and str(parser_label) != "":
                    page_label = str(parser_label)

            physical_page = page_index + 1
            drafts.append(
                _UnitDraft(
                    locator=SourceLocator(
                        locator_type=LocatorType.PAGE,
                        start=physical_page,
                        end=physical_page,
                        physical_page=physical_page,
                        page_label=page_label,
                    ),
                    section_path=section_paths[page_index],
                    raw_text=raw_text,
                    blocks=blocks,
                    joiner="",
                    has_images=has_images,
                    low_text_threshold=parse_config.pdf_low_text_characters,
                    metadata={
                        "width_points": round(width, 3),
                        "height_points": round(height, 3),
                        "rotation_degrees": int(page.rotation),
                        "image_count": image_count,
                        "explicit_page_labels": has_explicit_labels,
                    },
                )
            )
    finally:
        document.close()

    return _build_document(
        source_path,
        parser_name="pymupdf",
        drafts=drafts,
        config=parse_config,
        source_name=source_name,
        remove_repeated_marginalia=True,
        metadata={"page_count": len(drafts)},
    )


OCR_PARSER_NAME = "vision-ocr-v1"
OCR_PARSER_VERSION = 1


async def parse_scanned_pdf_document(
    path: str | Path,
    *,
    config: IngestionConfig | None = None,
    source_name: str | None = None,
    transcribe: Any,
    render_dpi: int = 120,
) -> IngestedDocument:
    """OCR a fully-scanned PDF by rendering pages and transcribing them.

    ``transcribe`` is an async callable ``(image_bytes: bytes) -> str`` supplied
    by the caller (the vision gateway).  Each page becomes one unit whose text
    is the transcription; equations are expected to arrive as LaTeX from the
    vision model.  Renders at ``render_dpi`` DPI as PNG for transport.
    """

    if not callable(transcribe):
        raise IngestionError("a transcribe callable is required for scanned PDFs")
    pymupdf = _load_pymupdf()
    source_path = Path(path)
    parse_config = config or IngestionConfig()
    drafts: list[_UnitDraft] = []
    document = None
    try:
        document = pymupdf.open(str(source_path))
    except Exception as error:
        raise IngestionError(f"could not open PDF {source_path.name!r}") from error

    try:
        page_count = int(document.page_count)
        section_paths = _toc_paths(document, page_count)
        has_explicit_labels = _pdf_has_explicit_page_labels(document)
        zoom = render_dpi / 72.0
        matrix = pymupdf.Matrix(zoom, zoom)
        for page_index in range(page_count):
            page = document.load_page(page_index)
            width = float(page.rect.width)
            height = float(page.rect.height)
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            image_bytes = pixmap.tobytes("png")

            transcription = ""
            transcription_error: str | None = None
            try:
                transcription = _sanitize_extracted_text(await transcribe(image_bytes)).strip()
            except Exception as error:
                # A single failed page must not abort the whole scan.  Preserve an
                # empty unit with a review flag so the gap is visible and the page
                # can be retried later without re-running the entire document.
                transcription_error = f"{type(error).__name__}: {error}"

            page_label: str | None = None
            if has_explicit_labels:
                try:
                    parser_label = page.get_label()
                except (AttributeError, RuntimeError, ValueError):
                    parser_label = None
                if parser_label is not None and str(parser_label) != "":
                    page_label = str(parser_label)

            physical_page = page_index + 1
            review_flags = {"RENDER_REQUIRED"}
            if transcription_error is not None:
                review_flags.add("OCR_FAILED")
            block = _DraftBlock(
                ordinal=0,
                exact_text=transcription,
                bbox=(0.0, 0.0, round(width, 3), round(height, 3)),
                kind="text",
                region="body",
                metadata={"source": "vision-ocr"},
            )
            drafts.append(
                _UnitDraft(
                    locator=SourceLocator(
                        locator_type=LocatorType.PAGE,
                        start=physical_page,
                        end=physical_page,
                        physical_page=physical_page,
                        page_label=page_label,
                    ),
                    section_path=section_paths[page_index],
                    raw_text=transcription,
                    blocks=[block],
                    joiner="",
                    has_images=True,
                    low_text_threshold=0,
                    force_review_flags=review_flags,
                    metadata={
                        "render_dpi": render_dpi,
                        "width_points": round(width, 3),
                        "height_points": round(height, 3),
                        "rotation_degrees": int(page.rotation),
                        "transcription_source": "vision-ocr",
                        "transcription_error": transcription_error,
                    },
                )
            )
    finally:
        if document is not None:
            document.close()

    return _build_document(
        source_path,
        parser_name=OCR_PARSER_NAME,
        drafts=drafts,
        config=parse_config,
        source_name=source_name,
        remove_repeated_marginalia=False,
        metadata={"page_count": len(drafts), "ocr_parser_version": OCR_PARSER_VERSION},
    )


_CHINESE_CHAPTER_RE = re.compile(r"^第[〇零一二三四五六七八九十百0-9]+章(?:\s|$)")
_NUMBERED_HEADING_RE = re.compile(r"^(?P<number>\d+(?:\.\d+)+)(?!\d)\s*")


def _docx_heading_level(text: str, style_name: str) -> int | None:
    style_match = re.search(r"(?:Heading|标题)\s*([1-9])", style_name, flags=re.IGNORECASE)
    if style_match:
        return int(style_match.group(1))
    if _CHINESE_CHAPTER_RE.match(text):
        return 1
    number_match = _NUMBERED_HEADING_RE.match(text)
    if number_match:
        return number_match.group("number").count(".") + 1
    return None


def _paragraph_has_drawing(paragraph: Any) -> bool:
    try:
        return bool(
            paragraph._p.xpath(
                ".//*[local-name()='drawing' or local-name()='pict' or local-name()='object']"
            )
        )
    except (AttributeError, TypeError, ValueError):
        return False


def parse_docx_document(
    path: str | Path,
    *,
    config: IngestionConfig | None = None,
    source_name: str | None = None,
) -> IngestedDocument:
    """Parse DOCX body paragraphs with source-native paragraph ordinals.

    DOCX does not reliably store rendered page boundaries, so this parser never
    emits a physical page number or a page label.
    """

    try:
        from docx import Document
    except ImportError as error:
        raise MissingParserDependency("DOCX ingestion requires python-docx") from error

    source_path = Path(path)
    parse_config = config or IngestionConfig()
    document = Document(str(source_path))
    drafts: list[_UnitDraft] = []
    root_title: str | None = None
    headings: dict[int, str] = {}

    for paragraph_ordinal, paragraph in enumerate(document.paragraphs, start=1):
        exact_text = _sanitize_extracted_text(paragraph.text)
        has_drawing = _paragraph_has_drawing(paragraph)
        if exact_text == "" and not has_drawing:
            continue
        style_name = paragraph.style.name if paragraph.style is not None else ""
        stripped = exact_text.strip()
        section_path: tuple[str, ...]

        is_title = style_name.casefold() == "title" or bool(
            re.fullmatch(r"《.+》教学大纲", stripped)
        )
        if is_title and stripped:
            root_title = exact_text
            headings.clear()
            section_path = (exact_text,)
        else:
            heading_level = _docx_heading_level(stripped, style_name)
            if heading_level is not None and stripped and not stripped.startswith(("*", "＊")):
                headings = {
                    level: heading for level, heading in headings.items() if level < heading_level
                }
                headings[heading_level] = exact_text
            ordered_headings = tuple(headings[level] for level in sorted(headings))
            section_path = ((root_title,) if root_title else ()) + ordered_headings

        flags: set[str] = set()
        if any(unicodedata.category(character) == "Co" for character in exact_text):
            flags.update(("PARSER_ANOMALY", "PRIVATE_USE_CHARACTER"))
        drafts.append(
            _UnitDraft(
                locator=SourceLocator(
                    locator_type=LocatorType.PARAGRAPH,
                    start=paragraph_ordinal,
                    end=paragraph_ordinal,
                ),
                section_path=section_path,
                raw_text=exact_text,
                blocks=[
                    _DraftBlock(
                        ordinal=0,
                        exact_text=exact_text,
                        kind="paragraph",
                        metadata={"style": style_name},
                    )
                ],
                has_images=has_drawing,
                low_text_threshold=1 if has_drawing else 0,
                force_review_flags=flags,
                metadata={"style": style_name, "contains_drawing": has_drawing},
            )
        )

    return _build_document(
        source_path,
        parser_name="python-docx",
        drafts=drafts,
        config=parse_config,
        source_name=source_name,
        metadata={"paragraph_count": len(document.paragraphs)},
    )


def _iter_pptx_shapes(shapes: Iterable[Any]) -> Iterator[Any]:
    for shape in shapes:
        yield shape
        try:
            nested_shapes = shape.shapes
        except (AttributeError, ValueError):
            continue
        yield from _iter_pptx_shapes(nested_shapes)


def _pptx_shape_text(shape: Any) -> tuple[str, str]:
    try:
        if shape.has_text_frame:
            return _sanitize_extracted_text(str(shape.text)), "text"
    except (AttributeError, ValueError):
        pass
    try:
        if shape.has_table:
            rows = []
            for row in shape.table.rows:
                rows.append("\t".join(_sanitize_extracted_text(cell.text) for cell in row.cells))
            return "\n".join(rows), "table"
    except (AttributeError, ValueError):
        pass
    return "", "graphic"


def parse_pptx_document(
    path: str | Path,
    *,
    config: IngestionConfig | None = None,
    source_name: str | None = None,
) -> IngestedDocument:
    """Parse presentation text without allowing chunks to cross slide boundaries."""

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as error:
        raise MissingParserDependency("PPTX ingestion requires python-pptx") from error

    source_path = Path(path)
    parse_config = config or IngestionConfig()
    presentation = Presentation(str(source_path))
    if presentation.slide_width is None or presentation.slide_height is None:
        raise IngestionError("PPTX presentation has no slide dimensions")
    slide_width = float(presentation.slide_width)
    slide_height = float(presentation.slide_height)
    drafts: list[_UnitDraft] = []

    picture_types = {
        value
        for value in (
            getattr(MSO_SHAPE_TYPE, "PICTURE", None),
            getattr(MSO_SHAPE_TYPE, "LINKED_PICTURE", None),
        )
        if value is not None
    }
    render_types = picture_types | {
        value
        for value in (
            getattr(MSO_SHAPE_TYPE, "CHART", None),
            getattr(MSO_SHAPE_TYPE, "DIAGRAM", None),
            getattr(MSO_SHAPE_TYPE, "MEDIA", None),
        )
        if value is not None
    }

    for slide_index, slide in enumerate(presentation.slides, start=1):
        flattened_shapes = list(_iter_pptx_shapes(slide.shapes))
        flattened_shapes.sort(
            key=lambda shape: (
                int(getattr(shape, "top", 0)),
                int(getattr(shape, "left", 0)),
            )
        )
        blocks: list[_DraftBlock] = []
        has_rendered_content = False
        image_count = 0
        for shape_ordinal, shape in enumerate(flattened_shapes):
            exact_text, kind = _pptx_shape_text(shape)
            shape_type = getattr(shape, "shape_type", None)
            if shape_type in picture_types:
                image_count += 1
            if shape_type in render_types:
                has_rendered_content = True
            if not exact_text and shape_type not in render_types:
                continue
            left = float(getattr(shape, "left", 0))
            top = float(getattr(shape, "top", 0))
            width = float(getattr(shape, "width", 0))
            height = float(getattr(shape, "height", 0))
            bottom = top + height
            if bottom <= slide_height * 0.12:
                region: Literal["top", "body", "bottom"] = "top"
            elif top >= slide_height * 0.88:
                region = "bottom"
            else:
                region = "body"
            blocks.append(
                _DraftBlock(
                    ordinal=shape_ordinal,
                    exact_text=exact_text,
                    bbox=(left, top, left + width, bottom),
                    kind="text" if exact_text else "image",
                    region=region,
                    metadata={
                        "shape_name": str(getattr(shape, "name", "")),
                        "shape_type": str(shape_type),
                        "content_kind": kind,
                        "bbox_unit": "EMU",
                    },
                )
            )

        title_text = ""
        try:
            title_shape = slide.shapes.title
            if title_shape is not None:
                title_text = str(title_shape.text)
        except (AttributeError, ValueError):
            pass
        section_path = (title_text,) if title_text else ()
        raw_text = "\n".join(
            block.exact_text for block in blocks if block.kind == "text" and block.exact_text
        )
        drafts.append(
            _UnitDraft(
                locator=SourceLocator(
                    locator_type=LocatorType.SLIDE,
                    start=slide_index,
                    end=slide_index,
                ),
                section_path=section_path,
                raw_text=raw_text,
                blocks=blocks,
                joiner="\n",
                has_images=image_count > 0 or has_rendered_content,
                low_text_threshold=parse_config.slide_low_text_characters,
                metadata={
                    "width_emu": int(slide_width),
                    "height_emu": int(slide_height),
                    "image_count": image_count,
                    "has_rendered_content": has_rendered_content,
                },
            )
        )

    return _build_document(
        source_path,
        parser_name="python-pptx",
        drafts=drafts,
        config=parse_config,
        source_name=source_name,
        remove_repeated_marginalia=True,
        metadata={"slide_count": len(drafts)},
    )


def _excel_cell_text(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, datetime):
        return value.isoformat(sep=" ")
    if isinstance(value, (date, time)):
        return value.isoformat()
    return _sanitize_extracted_text(str(value))


def parse_xlsx_document(
    path: str | Path,
    *,
    config: IngestionConfig | None = None,
    source_name: str | None = None,
) -> IngestedDocument:
    """Parse formulas and values verbatim, one source unit per non-empty sheet row."""

    try:
        from openpyxl import load_workbook  # type: ignore[import-untyped]
    except ImportError as error:
        raise MissingParserDependency("XLSX ingestion requires openpyxl") from error

    source_path = Path(path)
    parse_config = config or IngestionConfig()
    workbook = load_workbook(source_path, data_only=False, read_only=False)
    drafts: list[_UnitDraft] = []
    sheet_names: list[str] = []
    try:
        for worksheet in workbook.worksheets:
            sheet_names.append(worksheet.title)
            for row_number in range(1, worksheet.max_row + 1):
                last_nonempty_column = 0
                for column_number in range(1, worksheet.max_column + 1):
                    if worksheet.cell(row_number, column_number).value is not None:
                        last_nonempty_column = column_number
                if last_nonempty_column == 0:
                    continue

                blocks: list[_DraftBlock] = []
                values: list[str] = []
                formula_count = 0
                for column_number in range(1, last_nonempty_column + 1):
                    cell = worksheet.cell(row_number, column_number)
                    exact_text = _excel_cell_text(cell.value)
                    values.append(exact_text)
                    if cell.data_type == "f":
                        formula_count += 1
                    blocks.append(
                        _DraftBlock(
                            ordinal=column_number - 1,
                            exact_text=exact_text,
                            bbox=(
                                float(column_number),
                                float(row_number),
                                float(column_number),
                                float(row_number),
                            ),
                            kind="text",
                            metadata={
                                "cell": cell.coordinate,
                                "cell_data_type": str(cell.data_type),
                                "number_format": str(cell.number_format),
                            },
                        )
                    )
                row_text = "\t".join(values)
                flags = {"HIDDEN_ROW"} if worksheet.row_dimensions[row_number].hidden else set()
                drafts.append(
                    _UnitDraft(
                        locator=SourceLocator(
                            locator_type=LocatorType.SHEET_ROW,
                            start=row_number,
                            end=row_number,
                            sheet_name=worksheet.title,
                        ),
                        section_path=(worksheet.title,),
                        raw_text=row_text,
                        blocks=blocks,
                        joiner="\t",
                        force_review_flags=flags,
                        metadata={
                            "last_nonempty_column": last_nonempty_column,
                            "formula_count": formula_count,
                            "hidden": bool(worksheet.row_dimensions[row_number].hidden),
                        },
                    )
                )
    finally:
        workbook.close()

    return _build_document(
        source_path,
        parser_name="openpyxl",
        drafts=drafts,
        config=parse_config,
        source_name=source_name,
        metadata={
            "sheet_count": len(sheet_names),
            "sheet_names": json.dumps(sheet_names, ensure_ascii=False),
        },
    )


def _decode_text_file(path: Path) -> tuple[str, str, bool]:
    source_bytes = path.read_bytes()
    has_utf8_bom = source_bytes.startswith(b"\xef\xbb\xbf")
    attempts = ("utf-8-sig", "gb18030")
    for encoding in attempts:
        try:
            return _sanitize_extracted_text(source_bytes.decode(encoding)), encoding, has_utf8_bom
        except UnicodeDecodeError:
            continue
    raise IngestionError(f"could not decode text file {path.name!r} as UTF-8 or GB18030")


_MARKDOWN_HEADING_RE = re.compile(r"^(?P<marks>#{1,6})[ \t]+(?P<title>.*?)(?:\r?\n)?$")


def _text_drafts(text: str, *, markdown: bool) -> list[_UnitDraft]:
    lines = text.splitlines(keepends=True)
    if not lines:
        lines = [""]
    if not markdown:
        return [
            _UnitDraft(
                locator=SourceLocator(
                    locator_type=LocatorType.LINE,
                    start=1,
                    end=max(1, len(lines)),
                ),
                section_path=(),
                raw_text=text,
                blocks=[
                    _DraftBlock(ordinal=index, exact_text=line, kind="line")
                    for index, line in enumerate(lines)
                ],
            )
        ]

    drafts: list[_UnitDraft] = []
    headings: dict[int, str] = {}
    unit_lines: list[str] = []
    unit_start = 1
    unit_path: tuple[str, ...] = ()

    def flush(end_line: int) -> None:
        nonlocal unit_lines
        if not unit_lines:
            return
        raw_text = "".join(unit_lines)
        drafts.append(
            _UnitDraft(
                locator=SourceLocator(
                    locator_type=LocatorType.LINE,
                    start=unit_start,
                    end=max(unit_start, end_line),
                ),
                section_path=unit_path,
                raw_text=raw_text,
                blocks=[
                    _DraftBlock(ordinal=index, exact_text=line, kind="line")
                    for index, line in enumerate(unit_lines)
                ],
            )
        )
        unit_lines = []

    for line_number, line in enumerate(lines, start=1):
        heading_match = _MARKDOWN_HEADING_RE.match(line)
        if heading_match:
            flush(line_number - 1)
            level = len(heading_match.group("marks"))
            title = heading_match.group("title")
            headings = {
                old_level: value for old_level, value in headings.items() if old_level < level
            }
            headings[level] = title
            unit_start = line_number
            unit_path = tuple(headings[key] for key in sorted(headings))
        elif not unit_lines:
            unit_start = line_number
            unit_path = tuple(headings[key] for key in sorted(headings))
        unit_lines.append(line)
    flush(len(lines))
    return drafts


def parse_text_document(
    path: str | Path,
    *,
    config: IngestionConfig | None = None,
    source_name: str | None = None,
) -> IngestedDocument:
    """Parse UTF-8/GB18030 Markdown or plain text with exact line boundaries."""

    source_path = Path(path)
    if source_path.suffix.casefold() not in {".md", ".txt"}:
        raise UnsupportedDocumentError("parse_text_document accepts only .md and .txt")
    parse_config = config or IngestionConfig()
    text, encoding, has_bom = _decode_text_file(source_path)
    is_markdown = source_path.suffix.casefold() == ".md"
    return _build_document(
        source_path,
        parser_name="markdown" if is_markdown else "plain-text",
        drafts=_text_drafts(text, markdown=is_markdown),
        config=parse_config,
        source_name=source_name,
        metadata={"encoding": encoding, "utf8_bom": has_bom},
    )


def parse_document(
    path: str | Path,
    *,
    config: IngestionConfig | None = None,
    source_name: str | None = None,
) -> IngestedDocument:
    """Dispatch a supported source file to its deterministic parser."""

    source_path = Path(path)
    if not source_path.is_file():
        raise FileNotFoundError(source_path)
    parsers = {
        ".pdf": parse_pdf_document,
        ".docx": parse_docx_document,
        ".pptx": parse_pptx_document,
        ".xlsx": parse_xlsx_document,
        ".md": parse_text_document,
        ".txt": parse_text_document,
    }
    try:
        parser = parsers[source_path.suffix.casefold()]
    except KeyError as error:
        raise UnsupportedDocumentError(
            f"unsupported course-material suffix: {source_path.suffix or '<none>'}"
        ) from error
    return parser(source_path, config=config, source_name=source_name)


class SeedNodeType(StrEnum):
    CATEGORY = "Category"
    CONCEPT = "Concept"


class SeedRelationType(StrEnum):
    PART_OF = "PART_OF"
    PREREQUISITE_OF = "PREREQUISITE_OF"
    RELATED_TO = "RELATED_TO"


class HierarchySeedNode(BaseModel):
    """A reviewable node authored in the course knowledge-map spreadsheet."""

    model_config = ConfigDict(extra="forbid", frozen=True)

    id: str
    node_type: SeedNodeType
    verbatim_label: str
    canonical_label: str
    aliases: tuple[str, ...] = ()
    hierarchy_path: tuple[str, ...]
    parent_id: str | None = None
    tags: tuple[str, ...] = ()
    knowledge_category: str | None = None
    description: str | None = None
    source_chunk_id: str
    source_locator: SourceLocator
    source_row_text: str
    review_required: bool = False
    review_flags: tuple[str, ...] = ()


class HierarchySeedRelation(BaseModel):
    """A typed relation candidate; unresolved endpoints are retained for review."""

    model_config = ConfigDict(extra="forbid", frozen=True)

    id: str
    relation_type: SeedRelationType
    source_node_id: str | None
    target_node_id: str | None
    verbatim_source_label: str
    canonical_source_label: str
    verbatim_target_label: str
    canonical_target_label: str
    declared_on_node_id: str
    source_chunk_id: str
    source_locator: SourceLocator
    evidence_snippet: str
    review_required: bool = False
    review_flags: tuple[str, ...] = ()


class KnowledgeHierarchySeed(BaseModel):
    """Course-authored hierarchy/prerequisite seed; not teacher-approved graph data."""

    model_config = ConfigDict(extra="forbid", frozen=True)

    document: IngestedDocument
    nodes: tuple[HierarchySeedNode, ...]
    relations: tuple[HierarchySeedRelation, ...]
    diagnostics: tuple[str, ...] = ()


def canonicalize_course_label(label: str) -> str:
    """Normalize representation only; never translate or silently correct terminology."""

    normalized = unicodedata.normalize("NFKC", label)
    return re.sub(r"\s+", " ", normalized).strip()


def _label_key(label: str) -> str:
    return canonicalize_course_label(label).casefold()


def _authored_aliases(
    verbatim_label: str, canonical_label: str
) -> tuple[tuple[str, ...], set[str]]:
    aliases: list[str] = []
    flags: set[str] = set()
    if verbatim_label != canonical_label:
        aliases.append(verbatim_label)
    match = re.fullmatch(r"(?P<primary>[^()]+?)\((?P<alias>[^()]+)\)", canonical_label)
    if match:
        primary = match.group("primary").strip()
        alias = match.group("alias").strip()
        if primary and alias and not re.search(r"\d+\s*学时$", alias):
            aliases.extend((primary, alias))
            flags.add("AUTHORED_PARENTHETICAL_ALIAS")
    deduplicated: list[str] = []
    seen = {_label_key(canonical_label)}
    for alias in aliases:
        key = _label_key(alias)
        if key and key not in seen:
            seen.add(key)
            deduplicated.append(alias)
    return tuple(deduplicated), flags


@dataclass(slots=True)
class _SeedNodeDraft:
    id: str
    node_type: SeedNodeType
    verbatim_label: str
    canonical_label: str
    aliases: tuple[str, ...]
    hierarchy_path: tuple[str, ...]
    parent_id: str | None
    tags: tuple[str, ...]
    knowledge_category: str | None
    description: str | None
    source_chunk_id: str
    source_locator: SourceLocator
    source_row_text: str
    review_flags: set[str] = field(default_factory=set)


def _split_semicolon_cell(value: object) -> tuple[str, ...]:
    if value is None:
        return ()
    return tuple(part.strip() for part in re.split(r"[;；]", str(value)) if part.strip())


def _row_chunk_map(document: IngestedDocument) -> dict[tuple[str, int], SourceChunk]:
    result: dict[tuple[str, int], SourceChunk] = {}
    for chunk in document.chunks:
        locator = chunk.locator
        if locator.locator_type != LocatorType.SHEET_ROW or locator.sheet_name is None:
            continue
        key = (locator.sheet_name, int(locator.start))
        result.setdefault(key, chunk)
    return result


def _common_path_prefix(left: Sequence[str], right: Sequence[str]) -> int:
    score = 0
    for left_value, right_value in zip(left, right, strict=False):
        if _label_key(left_value) != _label_key(right_value):
            break
        score += 1
    return score


def _resolve_seed_reference(
    label: str,
    declared_node: HierarchySeedNode,
    nodes_by_id: Mapping[str, HierarchySeedNode],
    canonical_index: Mapping[str, Sequence[str]],
    alias_index: Mapping[str, Sequence[str]],
) -> tuple[str | None, set[str]]:
    key = _label_key(label)
    candidate_ids = list(canonical_index.get(key, ()))
    flags: set[str] = set()
    if not candidate_ids:
        candidate_ids = list(alias_index.get(key, ()))
        if candidate_ids:
            flags.add("ALIAS_REFERENCE_RESOLUTION")
    if not candidate_ids:
        return None, {"UNRESOLVED_REFERENCE"}
    if len(candidate_ids) == 1:
        return candidate_ids[0], flags

    declared_parent_path = declared_node.hierarchy_path[:-1]
    scored = [
        (
            _common_path_prefix(declared_parent_path, nodes_by_id[node_id].hierarchy_path[:-1]),
            node_id,
        )
        for node_id in candidate_ids
    ]
    best_score = max(score for score, _ in scored)
    best_ids = [node_id for score, node_id in scored if score == best_score]
    if best_score > 0 and len(best_ids) == 1:
        flags.add("CONTEXTUAL_REFERENCE_RESOLUTION")
        return best_ids[0], flags
    return None, {"AMBIGUOUS_REFERENCE"}


def parse_xlsx_hierarchy_seed(
    path: str | Path,
    *,
    config: IngestionConfig | None = None,
    source_name: str | None = None,
) -> KnowledgeHierarchySeed:
    """Parse the authored hierarchy, prerequisites, successors, and related nodes.

    The current course workbook uses columns B-H as hierarchy depth and I-K as
    prerequisite/successor/related references.  References that cannot be
    resolved exactly are preserved with a null endpoint and an explicit flag.
    """

    try:
        from openpyxl import load_workbook
    except ImportError as error:
        raise MissingParserDependency("XLSX seed ingestion requires openpyxl") from error

    source_path = Path(path)
    document = parse_xlsx_document(source_path, config=config, source_name=source_name)
    chunk_by_row = _row_chunk_map(document)
    workbook = load_workbook(source_path, data_only=False, read_only=False)
    node_drafts: list[_SeedNodeDraft] = []
    explicit_relations: list[tuple[_SeedNodeDraft, SeedRelationType, str, str]] = []
    diagnostics: list[str] = []

    try:
        for worksheet in workbook.worksheets:
            header_row: int | None = None
            for row_number in range(1, min(worksheet.max_row, 30) + 1):
                if str(worksheet.cell(row_number, 1).value or "").strip() == "节点类型*":
                    header_row = row_number
                    break
            if header_row is None:
                diagnostics.append(f"{worksheet.title}: missing 节点类型* header; sheet skipped")
                continue

            active_by_depth: dict[int, _SeedNodeDraft] = {}
            for row_number in range(header_row + 1, worksheet.max_row + 1):
                raw_type = worksheet.cell(row_number, 1).value
                label_cells = [
                    (depth, worksheet.cell(row_number, depth + 1).value)
                    for depth in range(1, 8)
                    if worksheet.cell(row_number, depth + 1).value not in (None, "")
                ]
                if raw_type in (None, "") and not label_cells:
                    continue
                type_text = str(raw_type or "").strip()
                node_type = {
                    "分类": SeedNodeType.CATEGORY,
                    "知识点": SeedNodeType.CONCEPT,
                }.get(type_text)
                if node_type is None:
                    diagnostics.append(
                        f"{worksheet.title}!{row_number}: unknown node type {type_text!r}"
                    )
                    continue
                if not label_cells:
                    diagnostics.append(f"{worksheet.title}!{row_number}: node row has no label")
                    continue

                row_chunk = chunk_by_row[(worksheet.title, row_number)]
                row_flags = {"MULTIPLE_LABEL_CELLS"} if len(label_cells) > 1 else set()
                created_this_row: list[_SeedNodeDraft] = []
                for depth, raw_label in label_cells:
                    verbatim_label = str(raw_label)
                    canonical_label = canonicalize_course_label(verbatim_label)
                    aliases, alias_flags = _authored_aliases(verbatim_label, canonical_label)
                    flags = set(row_flags) | alias_flags
                    if any(unicodedata.category(character) == "Co" for character in verbatim_label):
                        flags.add("PRIVATE_USE_CHARACTER")

                    active_by_depth = {
                        old_depth: node
                        for old_depth, node in active_by_depth.items()
                        if old_depth < depth
                    }
                    parent_depths = [
                        old_depth for old_depth in active_by_depth if old_depth < depth
                    ]
                    parent = active_by_depth[max(parent_depths)] if parent_depths else None
                    if depth > 1 and parent is None:
                        flags.add("MISSING_PARENT")
                    hierarchy_path = (
                        (*parent.hierarchy_path, verbatim_label)
                        if parent is not None
                        else (verbatim_label,)
                    )
                    node_id = _stable_id(
                        "seednode",
                        document.document_version_id,
                        worksheet.title,
                        row_number,
                        depth,
                        verbatim_label,
                    )
                    tags = _split_semicolon_cell(worksheet.cell(row_number, 12).value)
                    knowledge_category_value = worksheet.cell(row_number, 13).value
                    description_value = worksheet.cell(row_number, 14).value
                    node = _SeedNodeDraft(
                        id=node_id,
                        node_type=node_type,
                        verbatim_label=verbatim_label,
                        canonical_label=canonical_label,
                        aliases=aliases,
                        hierarchy_path=hierarchy_path,
                        parent_id=parent.id if parent is not None else None,
                        tags=tags,
                        knowledge_category=(
                            str(knowledge_category_value)
                            if knowledge_category_value not in (None, "")
                            else None
                        ),
                        description=(
                            str(description_value) if description_value not in (None, "") else None
                        ),
                        source_chunk_id=row_chunk.id,
                        source_locator=row_chunk.locator,
                        source_row_text=row_chunk.exact_text,
                        review_flags=flags,
                    )
                    node_drafts.append(node)
                    created_this_row.append(node)
                    active_by_depth[depth] = node

                declared_node = created_this_row[-1]
                for column, relation_type, declaration_kind in (
                    (9, SeedRelationType.PREREQUISITE_OF, "prerequisite"),
                    (10, SeedRelationType.PREREQUISITE_OF, "successor"),
                    (11, SeedRelationType.RELATED_TO, "related"),
                ):
                    raw_relation_cell = worksheet.cell(row_number, column).value
                    for referenced_label in _split_semicolon_cell(raw_relation_cell):
                        explicit_relations.append(
                            (declared_node, relation_type, declaration_kind, referenced_label)
                        )
    finally:
        workbook.close()

    duplicate_index: dict[str, list[_SeedNodeDraft]] = defaultdict(list)
    for draft_node in node_drafts:
        duplicate_index[_label_key(draft_node.canonical_label)].append(draft_node)
    for duplicates in duplicate_index.values():
        if len(duplicates) > 1:
            for duplicate_node in duplicates:
                duplicate_node.review_flags.add("DUPLICATE_CANONICAL_LABEL")

    nodes = tuple(
        HierarchySeedNode(
            id=draft_node.id,
            node_type=draft_node.node_type,
            verbatim_label=draft_node.verbatim_label,
            canonical_label=draft_node.canonical_label,
            aliases=draft_node.aliases,
            hierarchy_path=draft_node.hierarchy_path,
            parent_id=draft_node.parent_id,
            tags=draft_node.tags,
            knowledge_category=draft_node.knowledge_category,
            description=draft_node.description,
            source_chunk_id=draft_node.source_chunk_id,
            source_locator=draft_node.source_locator,
            source_row_text=draft_node.source_row_text,
            review_required=bool(draft_node.review_flags),
            review_flags=tuple(sorted(draft_node.review_flags)),
        )
        for draft_node in node_drafts
    )
    nodes_by_id = {final_node.id: final_node for final_node in nodes}
    canonical_index: dict[str, list[str]] = defaultdict(list)
    alias_index: dict[str, list[str]] = defaultdict(list)
    for final_node in nodes:
        canonical_index[_label_key(final_node.canonical_label)].append(final_node.id)
        for alias in final_node.aliases:
            alias_index[_label_key(alias)].append(final_node.id)

    relations: list[HierarchySeedRelation] = []
    for final_node in nodes:
        if final_node.parent_id is None:
            continue
        parent_node = nodes_by_id[final_node.parent_id]
        relation_flags = set(final_node.review_flags) & {
            "MISSING_PARENT",
            "MULTIPLE_LABEL_CELLS",
        }
        relations.append(
            HierarchySeedRelation(
                id=_stable_id(
                    "seedrel",
                    final_node.id,
                    SeedRelationType.PART_OF.value,
                    parent_node.id,
                ),
                relation_type=SeedRelationType.PART_OF,
                source_node_id=final_node.id,
                target_node_id=parent_node.id,
                verbatim_source_label=final_node.verbatim_label,
                canonical_source_label=final_node.canonical_label,
                verbatim_target_label=parent_node.verbatim_label,
                canonical_target_label=parent_node.canonical_label,
                declared_on_node_id=final_node.id,
                source_chunk_id=final_node.source_chunk_id,
                source_locator=final_node.source_locator,
                evidence_snippet=final_node.source_row_text,
                review_required=bool(relation_flags),
                review_flags=tuple(sorted(relation_flags)),
            )
        )

    for declared_draft, relation_type, declaration_kind, referenced_label in explicit_relations:
        declared_final_node = nodes_by_id[declared_draft.id]
        resolved_id, resolution_flags = _resolve_seed_reference(
            referenced_label,
            declared_final_node,
            nodes_by_id,
            canonical_index,
            alias_index,
        )
        source_node_id: str | None
        target_node_id: str | None
        if declaration_kind == "prerequisite":
            source_node_id = resolved_id
            target_node_id = declared_final_node.id
            verbatim_source = referenced_label
            canonical_source = canonicalize_course_label(referenced_label)
            verbatim_target = declared_final_node.verbatim_label
            canonical_target = declared_final_node.canonical_label
        else:
            source_node_id = declared_final_node.id
            target_node_id = resolved_id
            verbatim_source = declared_final_node.verbatim_label
            canonical_source = declared_final_node.canonical_label
            verbatim_target = referenced_label
            canonical_target = canonicalize_course_label(referenced_label)
        if resolved_id == declared_final_node.id:
            resolution_flags.add("SELF_RELATION")
        relations.append(
            HierarchySeedRelation(
                id=_stable_id(
                    "seedrel",
                    document.document_version_id,
                    declared_final_node.id,
                    relation_type.value,
                    declaration_kind,
                    referenced_label,
                ),
                relation_type=relation_type,
                source_node_id=source_node_id,
                target_node_id=target_node_id,
                verbatim_source_label=verbatim_source,
                canonical_source_label=canonical_source,
                verbatim_target_label=verbatim_target,
                canonical_target_label=canonical_target,
                declared_on_node_id=declared_final_node.id,
                source_chunk_id=declared_final_node.source_chunk_id,
                source_locator=declared_final_node.source_locator,
                evidence_snippet=declared_final_node.source_row_text,
                review_required=bool(resolution_flags),
                review_flags=tuple(sorted(resolution_flags)),
            )
        )

    return KnowledgeHierarchySeed(
        document=document,
        nodes=nodes,
        relations=tuple(relations),
        diagnostics=tuple(diagnostics),
    )


class OutlineEntryType(StrEnum):
    COURSE_TITLE = "COURSE_TITLE"
    TERM = "TERM"
    CHAPTER = "CHAPTER"
    SECTION = "SECTION"
    TOPIC = "TOPIC"
    REFERENCE = "REFERENCE"


class CourseOutlineEntry(BaseModel):
    """One exact paragraph from the authored course outline."""

    model_config = ConfigDict(extra="forbid", frozen=True)

    id: str
    entry_type: OutlineEntryType
    verbatim_label: str
    canonical_label: str
    outline_number: str | None = None
    depth: int = Field(ge=0)
    hours: int | None = Field(default=None, ge=0)
    parent_id: str | None = None
    section_path: tuple[str, ...] = ()
    source_chunk_id: str
    source_locator: SourceLocator
    review_required: bool = False
    review_flags: tuple[str, ...] = ()

    @model_validator(mode="after")
    def forbid_fake_pages(self) -> CourseOutlineEntry:
        if self.source_locator.locator_type != LocatorType.PARAGRAPH:
            raise ValueError("DOCX outline entries must use paragraph locators")
        if (
            self.source_locator.physical_page is not None
            or self.source_locator.page_label is not None
        ):
            raise ValueError("DOCX outline entries cannot claim rendered page metadata")
        return self


class CourseOutlineRelation(BaseModel):
    model_config = ConfigDict(extra="forbid", frozen=True)

    id: str
    relation_type: Literal["PART_OF"] = "PART_OF"
    child_id: str
    parent_id: str
    source_chunk_id: str


class CourseOutlineSeed(BaseModel):
    """Parsed teaching organization from the DOCX outline, pending review."""

    model_config = ConfigDict(extra="forbid", frozen=True)

    document: IngestedDocument
    academic_year: int | None = None
    term: str | None = None
    entries: tuple[CourseOutlineEntry, ...]
    relations: tuple[CourseOutlineRelation, ...]
    diagnostics: tuple[str, ...] = ()


_OUTLINE_CHAPTER_RE = re.compile(
    r"^(?P<number>第[〇零一二三四五六七八九十百0-9]+章)\s*(?P<title>.*?)"
    r"(?:\s*[（(](?P<hours>\d+)\s*学时[）)])?\s*$"
)
_OUTLINE_SECTION_RE = re.compile(r"^(?P<number>\d+(?:\.\d+)+)(?!\d)\s*(?P<title>.*)$")
_OUTLINE_TERM_RE = re.compile(r"^-+\s*(?P<year>20\d{2})\s*[，,]\s*(?P<term>春|夏|秋|冬)\s*-+$")


@dataclass(slots=True)
class _OutlineEntryDraft:
    id: str
    entry_type: OutlineEntryType
    verbatim_label: str
    canonical_label: str
    outline_number: str | None
    depth: int
    hours: int | None
    parent_id: str | None
    section_path: tuple[str, ...]
    source_chunk_id: str
    source_locator: SourceLocator
    review_flags: set[str] = field(default_factory=set)


def _paragraph_chunk_map(document: IngestedDocument) -> dict[int, SourceChunk]:
    result: dict[int, SourceChunk] = {}
    for chunk in document.chunks:
        if chunk.locator.locator_type != LocatorType.PARAGRAPH:
            continue
        result.setdefault(int(chunk.locator.start), chunk)
    return result


def parse_docx_outline_seed(
    path: str | Path,
    *,
    config: IngestionConfig | None = None,
    source_name: str | None = None,
) -> CourseOutlineSeed:
    """Parse the 2026-style course outline without manufacturing page numbers.

    Source typos and duplicate section numbers are retained verbatim and marked
    for review.  This is intentionally not a curriculum autocorrector.
    """

    document = parse_docx_document(path, config=config, source_name=source_name)
    chunk_by_paragraph = _paragraph_chunk_map(document)
    entry_drafts: list[_OutlineEntryDraft] = []
    diagnostics: list[str] = []
    title_id: str | None = None
    current_chapter_id: str | None = None
    entry_by_number: dict[str, _OutlineEntryDraft] = {}
    entries_by_number: dict[str, list[_OutlineEntryDraft]] = defaultdict(list)
    current_section_by_depth: dict[int, _OutlineEntryDraft] = {}
    academic_year: int | None = None
    term: str | None = None

    for unit in document.units:
        paragraph_ordinal = int(unit.locator.start)
        verbatim = unit.raw_text
        stripped = verbatim.strip()
        if not stripped:
            continue
        chunk = chunk_by_paragraph[paragraph_ordinal]
        flags: set[str] = set()
        entry_type: OutlineEntryType
        outline_number: str | None
        depth: int
        hours: int | None
        parent_id: str | None
        section_path: tuple[str, ...]
        if any(unicodedata.category(character) == "Co" for character in verbatim):
            flags.add("PRIVATE_USE_CHARACTER")

        chapter_match = _OUTLINE_CHAPTER_RE.match(stripped)
        section_match = _OUTLINE_SECTION_RE.match(stripped)
        term_match = _OUTLINE_TERM_RE.match(stripped)
        if re.fullmatch(r"《.+》教学大纲", stripped):
            entry_type = OutlineEntryType.COURSE_TITLE
            outline_number = None
            depth = 0
            hours = None
            parent_id = None
            section_path = (verbatim,)
        elif term_match:
            entry_type = OutlineEntryType.TERM
            outline_number = None
            depth = 1
            hours = None
            parent_id = title_id
            section_path = ()
            academic_year = int(term_match.group("year"))
            term = term_match.group("term")
        elif chapter_match:
            entry_type = OutlineEntryType.CHAPTER
            outline_number = chapter_match.group("number")
            depth = 1
            hours_value = chapter_match.group("hours")
            hours = int(hours_value) if hours_value is not None else None
            parent_id = title_id
            section_path = (verbatim,)
        elif section_match:
            entry_type = OutlineEntryType.SECTION
            outline_number = section_match.group("number")
            depth = outline_number.count(".") + 1
            hours = None
            parent_number = outline_number.rsplit(".", 1)[0]
            parent_entry = entry_by_number.get(parent_number)
            if parent_entry is not None:
                parent_id = parent_entry.id
                section_path = (*parent_entry.section_path, verbatim)
            else:
                parent_id = current_chapter_id
                chapter_entry = next(
                    (entry for entry in reversed(entry_drafts) if entry.id == current_chapter_id),
                    None,
                )
                section_path = (
                    (*chapter_entry.section_path, verbatim)
                    if chapter_entry is not None
                    else (verbatim,)
                )
                if depth > 2:
                    flags.add("MISSING_NUMBERED_PARENT")
        elif stripped.startswith(("*", "＊")):
            entry_type = OutlineEntryType.REFERENCE
            outline_number = None
            depth = 2
            hours = None
            parent_id = current_chapter_id or title_id
            chapter_entry = next(
                (entry for entry in reversed(entry_drafts) if entry.id == current_chapter_id),
                None,
            )
            section_path = chapter_entry.section_path if chapter_entry is not None else ()
        else:
            entry_type = OutlineEntryType.TOPIC
            outline_number = None
            parent_section = (
                current_section_by_depth[max(current_section_by_depth)]
                if current_section_by_depth
                else None
            )
            parent_id = parent_section.id if parent_section is not None else current_chapter_id
            depth = (parent_section.depth + 1) if parent_section is not None else 2
            hours = None
            section_path = parent_section.section_path if parent_section is not None else ()
            flags.add("UNNUMBERED_TOPIC")

        entry_id = _stable_id(
            "outline",
            document.document_version_id,
            paragraph_ordinal,
            verbatim,
        )
        entry = _OutlineEntryDraft(
            id=entry_id,
            entry_type=entry_type,
            verbatim_label=verbatim,
            canonical_label=canonicalize_course_label(verbatim),
            outline_number=outline_number,
            depth=depth,
            hours=hours,
            parent_id=parent_id,
            section_path=section_path,
            source_chunk_id=chunk.id,
            source_locator=chunk.locator,
            review_flags=flags,
        )
        entry_drafts.append(entry)

        if entry_type == OutlineEntryType.COURSE_TITLE:
            title_id = entry.id
        elif entry_type == OutlineEntryType.CHAPTER:
            current_chapter_id = entry.id
            current_section_by_depth.clear()
            if outline_number is not None:
                entry_by_number[outline_number] = entry
                entries_by_number[outline_number].append(entry)
        elif entry_type == OutlineEntryType.SECTION and outline_number is not None:
            current_section_by_depth = {
                old_depth: old_entry
                for old_depth, old_entry in current_section_by_depth.items()
                if old_depth < depth
            }
            current_section_by_depth[depth] = entry
            entry_by_number[outline_number] = entry
            entries_by_number[outline_number].append(entry)

    for outline_number, duplicate_entries in entries_by_number.items():
        if len(duplicate_entries) <= 1:
            continue
        diagnostics.append(
            f"duplicate outline number {outline_number!r} at paragraphs "
            + ", ".join(str(entry.source_locator.start) for entry in duplicate_entries)
        )
        for entry in duplicate_entries:
            entry.review_flags.add("DUPLICATE_OUTLINE_NUMBER")

    entries = tuple(
        CourseOutlineEntry(
            id=entry.id,
            entry_type=entry.entry_type,
            verbatim_label=entry.verbatim_label,
            canonical_label=entry.canonical_label,
            outline_number=entry.outline_number,
            depth=entry.depth,
            hours=entry.hours,
            parent_id=entry.parent_id,
            section_path=entry.section_path,
            source_chunk_id=entry.source_chunk_id,
            source_locator=entry.source_locator,
            review_required=bool(entry.review_flags),
            review_flags=tuple(sorted(entry.review_flags)),
        )
        for entry in entry_drafts
    )
    relations = tuple(
        CourseOutlineRelation(
            id=_stable_id("outlinerel", entry.id, entry.parent_id),
            child_id=entry.id,
            parent_id=entry.parent_id,
            source_chunk_id=entry.source_chunk_id,
        )
        for entry in entries
        if entry.parent_id is not None
    )
    return CourseOutlineSeed(
        document=document,
        academic_year=academic_year,
        term=term,
        entries=entries,
        relations=relations,
        diagnostics=tuple(diagnostics),
    )


def parse_2026_course_outline(
    path: str | Path,
    *,
    config: IngestionConfig | None = None,
    source_name: str | None = None,
) -> CourseOutlineSeed:
    """Compatibility entry point naming the course's current outline edition."""

    return parse_docx_outline_seed(path, config=config, source_name=source_name)


__all__ = [
    "CourseOutlineEntry",
    "CourseOutlineRelation",
    "CourseOutlineSeed",
    "HierarchySeedNode",
    "HierarchySeedRelation",
    "IngestedDocument",
    "IngestionConfig",
    "IngestionError",
    "IngestionStatus",
    "KnowledgeHierarchySeed",
    "LocatorType",
    "MissingParserDependency",
    "OutlineEntryType",
    "SeedNodeType",
    "SeedRelationType",
    "SourceChunk",
    "SourceLocator",
    "SourceUnit",
    "TextBlock",
    "UnsupportedDocumentError",
    "canonicalize_course_label",
    "parse_2026_course_outline",
    "parse_document",
    "parse_docx_document",
    "parse_docx_outline_seed",
    "parse_pdf_document",
    "parse_pptx_document",
    "parse_text_document",
    "parse_xlsx_document",
    "parse_xlsx_hierarchy_seed",
    "sha256_bytes",
    "sha256_text",
]
